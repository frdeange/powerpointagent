"""
PPTX MCP Server — FastMCP 3.x
Exposes 7 tools for building and analyzing PowerPoint presentations.
Hosted on Azure Container Apps with external ingress.
Azure AI Foundry calls this server directly via AzureAIClient.get_mcp_tool().
"""

import os
import io
import json
import uuid
import logging
from pathlib import Path
from typing import Any

from fastmcp import FastMCP
from fastmcp.server.middleware import RateLimitMiddleware
from starlette.requests import Request
from starlette.responses import JSONResponse

from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
from datetime import datetime, timedelta, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ── FastMCP setup ────────────────────────────────────────────────────────────
mcp = FastMCP(
    name="pptx-mcp-server",
    instructions=(
        "Provides tools to create, modify, and analyze PowerPoint presentations. "
        "All presentations are stored in Azure Blob Storage. "
        "Always call create_presentation first, then add_slide for each slide, "
        "and finally export_presentation to get the download URL."
    ),
    stateless_http=True,
)

# ── Azure Storage client ─────────────────────────────────────────────────────
_blob_client: BlobServiceClient | None = None


def get_blob_client() -> BlobServiceClient:
    global _blob_client
    if _blob_client is None:
        conn_str = os.environ["AZURE_STORAGE_CONNECTION_STRING"]
        _blob_client = BlobServiceClient.from_connection_string(conn_str)
    return _blob_client


CONTAINER_GENERATED = os.environ.get("BLOB_CONTAINER_GENERATED", "generated")
CONTAINER_TEMPLATES = os.environ.get("BLOB_CONTAINER_TEMPLATES", "templates")
CONTAINER_UPLOADS = os.environ.get("BLOB_CONTAINER_UPLOADS", "uploads")
STORAGE_ACCOUNT = os.environ.get("AZURE_STORAGE_ACCOUNT_NAME", "")


def _generate_sas_url(container: str, blob_name: str, expiry_hours: int = 24) -> str:
    """Generate a SAS URL for blob download."""
    from azure.storage.blob import generate_blob_sas, BlobSasPermissions

    client = get_blob_client()
    account_name = client.account_name
    account_key = client.credential.account_key

    sas_token = generate_blob_sas(
        account_name=account_name,
        container_name=container,
        blob_name=blob_name,
        account_key=account_key,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.now(timezone.utc) + timedelta(hours=expiry_hours),
    )
    return f"https://{account_name}.blob.core.windows.net/{container}/{blob_name}?{sas_token}"


def _download_blob(container: str, blob_name: str) -> bytes:
    """Download a blob and return its bytes."""
    client = get_blob_client()
    blob = client.get_blob_client(container=container, blob=blob_name)
    return blob.download_blob().readall()


def _upload_blob(
    container: str,
    blob_name: str,
    data: bytes,
    content_type: str = "application/octet-stream",
) -> None:
    """Upload bytes to a blob."""
    client = get_blob_client()
    blob = client.get_blob_client(container=container, blob=blob_name)
    blob.upload_blob(
        data, overwrite=True, content_settings={"content_type": content_type}
    )


# ── Tool 1: create_presentation ──────────────────────────────────────────────
@mcp.tool()
def create_presentation(
    title: str,
    subtitle: str = "",
    template_name: str = "default",
    presentation_id: str = "",
) -> dict[str, Any]:
    """
    Initialize a new PowerPoint presentation in Azure Blob Storage.

    Args:
        title: Title of the presentation.
        subtitle: Optional subtitle for the title slide.
        template_name: Name of the design template to apply (use list_templates to see available ones).
        presentation_id: Optional custom ID. A UUID is generated if not provided.

    Returns:
        dict with presentation_id, blob_name, and status.
    """
    pptx_id = presentation_id or str(uuid.uuid4())
    blob_name = f"{pptx_id}.pptx"

    # Try to load template from blob storage, else use blank
    prs = Presentation()
    try:
        tmpl_bytes = _download_blob(CONTAINER_TEMPLATES, f"{template_name}.pptx")
        prs = Presentation(io.BytesIO(tmpl_bytes))
        logger.info("Loaded template: %s", template_name)
    except Exception:
        logger.warning(
            "Template '%s' not found, using blank presentation.", template_name
        )

    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:  # subtitle
            ph.text = subtitle

    # Save to blob storage
    buf = io.BytesIO()
    prs.save(buf)
    _upload_blob(CONTAINER_GENERATED, blob_name, buf.getvalue())

    logger.info("Created presentation: %s", pptx_id)
    return {
        "presentation_id": pptx_id,
        "blob_name": blob_name,
        "slide_count": len(prs.slides),
        "status": "created",
        "template_used": template_name,
    }


# ── Tool 2: add_slide ────────────────────────────────────────────────────────
@mcp.tool()
def add_slide(
    presentation_id: str,
    slide_title: str,
    content: list[str],
    layout: str = "content",
    speaker_notes: str = "",
    image_url: str = "",
) -> dict[str, Any]:
    """
    Add a slide to an existing presentation.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_title: Title of the slide.
        content: List of bullet points or content lines.
        layout: Slide layout — 'content', 'two_column', 'image_only', 'blank', 'section_header'.
        speaker_notes: Speaker notes text for this slide.
        image_url: Optional URL of an image to embed (blob SAS URL or public URL).

    Returns:
        dict with slide_index and status.
    """
    blob_name = f"{presentation_id}.pptx"
    pptx_bytes = _download_blob(CONTAINER_GENERATED, blob_name)
    prs = Presentation(io.BytesIO(pptx_bytes))

    layout_map = {
        "content": 1,
        "two_column": 3,
        "image_only": 5,
        "blank": 6,
        "section_header": 2,
        "title_only": 5,
    }
    layout_idx = layout_map.get(layout, 1)
    # Clamp to available layouts
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_title

    # Set content placeholder
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and content:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(content):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = bullet
                para.level = 0

    # Add image if provided
    if image_url:
        try:
            import urllib.request

            with urllib.request.urlopen(image_url, timeout=20) as resp:
                img_bytes = resp.read()
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                img_stream,
                left=Inches(7.5),
                top=Inches(1.5),
                width=Inches(5.0),
                height=Inches(4.5),
            )
        except Exception as e:
            logger.warning("Could not embed image: %s", e)

    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes

    # Save back
    buf = io.BytesIO()
    prs.save(buf)
    _upload_blob(CONTAINER_GENERATED, blob_name, buf.getvalue())

    slide_index = len(prs.slides) - 1
    logger.info("Added slide %d to presentation %s", slide_index, presentation_id)
    return {
        "presentation_id": presentation_id,
        "slide_index": slide_index,
        "slide_title": slide_title,
        "layout": layout,
        "status": "slide_added",
    }


# ── Tool 3: apply_template ───────────────────────────────────────────────────
@mcp.tool()
def apply_template(
    presentation_id: str,
    template_name: str,
    primary_color_hex: str = "",
    font_name: str = "",
) -> dict[str, Any]:
    """
    Apply a design template or color/font overrides to an existing presentation.

    Args:
        presentation_id: ID of the presentation to modify.
        template_name: Name of the template (e.g. 'corporate_blue', 'minimal_white').
        primary_color_hex: Optional hex color to override primary accent (e.g. '#0078D4').
        font_name: Optional font name override (e.g. 'Segoe UI', 'Calibri').

    Returns:
        dict with status and changes applied.
    """
    blob_name = f"{presentation_id}.pptx"
    pptx_bytes = _download_blob(CONTAINER_GENERATED, blob_name)
    prs = Presentation(io.BytesIO(pptx_bytes))

    changes: list[str] = []

    # Apply color override to all slides
    if primary_color_hex:
        hex_clean = primary_color_hex.lstrip("#")
        try:
            r, g, b = (
                int(hex_clean[0:2], 16),
                int(hex_clean[2:4], 16),
                int(hex_clean[4:6], 16),
            )
            color = RGBColor(r, g, b)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color.type is not None:
                                    run.font.color.rgb = color
            changes.append(f"primary_color={primary_color_hex}")
        except Exception as e:
            logger.warning("Could not apply color: %s", e)

    # Apply font override
    if font_name:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = font_name
        changes.append(f"font={font_name}")

    buf = io.BytesIO()
    prs.save(buf)
    _upload_blob(CONTAINER_GENERATED, blob_name, buf.getvalue())

    return {
        "presentation_id": presentation_id,
        "template_applied": template_name,
        "changes": changes,
        "status": "template_applied",
    }


# ── Tool 4: add_image_to_slide ───────────────────────────────────────────────
@mcp.tool()
def add_image_to_slide(
    presentation_id: str,
    slide_index: int,
    image_url: str,
    left_inches: float = 1.0,
    top_inches: float = 2.0,
    width_inches: float = 5.0,
    height_inches: float = 3.5,
) -> dict[str, Any]:
    """
    Insert an image into a specific slide.

    Args:
        presentation_id: ID of the presentation.
        slide_index: 0-based index of the slide to modify.
        image_url: SAS URL or public URL of the image.
        left_inches: Distance from left edge (inches).
        top_inches: Distance from top edge (inches).
        width_inches: Image width (inches).
        height_inches: Image height (inches).

    Returns:
        dict with slide_index, shape_name, and status.
    """
    import urllib.request

    blob_name = f"{presentation_id}.pptx"
    pptx_bytes = _download_blob(CONTAINER_GENERATED, blob_name)
    prs = Presentation(io.BytesIO(pptx_bytes))

    if slide_index >= len(prs.slides):
        return {
            "error": f"slide_index {slide_index} out of range (total: {len(prs.slides)})"
        }

    slide = prs.slides[slide_index]
    with urllib.request.urlopen(image_url, timeout=20) as resp:
        img_bytes = resp.read()

    pic = slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        left=Inches(left_inches),
        top=Inches(top_inches),
        width=Inches(width_inches),
        height=Inches(height_inches),
    )

    buf = io.BytesIO()
    prs.save(buf)
    _upload_blob(CONTAINER_GENERATED, blob_name, buf.getvalue())

    return {
        "presentation_id": presentation_id,
        "slide_index": slide_index,
        "shape_name": pic.name,
        "status": "image_added",
    }


# ── Tool 5: export_presentation ──────────────────────────────────────────────
@mcp.tool()
def export_presentation(
    presentation_id: str,
    expiry_hours: int = 48,
) -> dict[str, Any]:
    """
    Finalize the presentation and generate a SAS download URL.

    Args:
        presentation_id: ID of the presentation to export.
        expiry_hours: How long the download URL remains valid (default 48h).

    Returns:
        dict with download_url, slide_count, file_size_kb, and expiry info.
    """
    blob_name = f"{presentation_id}.pptx"
    pptx_bytes = _download_blob(CONTAINER_GENERATED, blob_name)
    prs = Presentation(io.BytesIO(pptx_bytes))

    download_url = _generate_sas_url(CONTAINER_GENERATED, blob_name, expiry_hours)
    file_size_kb = round(len(pptx_bytes) / 1024, 1)
    expiry_time = (
        datetime.now(timezone.utc) + timedelta(hours=expiry_hours)
    ).isoformat()

    logger.info(
        "Exported presentation %s (%d slides, %.1f KB)",
        presentation_id,
        len(prs.slides),
        file_size_kb,
    )
    return {
        "presentation_id": presentation_id,
        "download_url": download_url,
        "slide_count": len(prs.slides),
        "file_size_kb": file_size_kb,
        "expires_at": expiry_time,
        "status": "exported",
    }


# ── Tool 6: list_templates ───────────────────────────────────────────────────
@mcp.tool()
def list_templates() -> dict[str, Any]:
    """
    List all available design templates in blob storage.

    Returns:
        dict with list of template names and descriptions.
    """
    client = get_blob_client()
    container = client.get_container_client(CONTAINER_TEMPLATES)

    templates = []
    try:
        for blob in container.list_blobs():
            if blob.name.endswith(".pptx"):
                name = blob.name.replace(".pptx", "")
                templates.append(
                    {"name": name, "size_kb": round((blob.size or 0) / 1024, 1)}
                )
    except Exception as e:
        logger.warning("Could not list templates: %s", e)

    # Always include built-in default
    if not any(t["name"] == "default" for t in templates):
        templates.insert(
            0,
            {"name": "default", "size_kb": 0, "description": "Blank 16:9 presentation"},
        )

    return {"templates": templates, "count": len(templates)}


# ── Tool 7: analyze_pptx_document ────────────────────────────────────────────
@mcp.tool()
def analyze_pptx_document(
    blob_name: str,
    container: str = "",
    extract_images: bool = False,
) -> dict[str, Any]:
    """
    Analyze an uploaded PowerPoint file to extract design specifications and content outline.
    Used by the DocumentAnalyzer agent to feed uploaded presentations into the generation pipeline.

    Args:
        blob_name: Name of the uploaded PPTX blob.
        container: Blob container (defaults to uploads container).
        extract_images: Whether to extract and enumerate image shapes.

    Returns:
        dict with design_spec (colors, fonts, layouts) and content_outline (slide titles, bullets).
    """
    source_container = container or CONTAINER_UPLOADS
    pptx_bytes = _download_blob(source_container, blob_name)
    prs = Presentation(io.BytesIO(pptx_bytes))

    slide_width_in = round(prs.slide_width / 914400, 2)  # EMU → inches
    slide_height_in = round(prs.slide_height / 914400, 2)

    # ── Extract content outline ──────────────────────────────────────────────
    content_outline: list[dict] = []
    fonts_found: set[str] = set()
    colors_found: set[str] = set()
    layouts_used: set[str] = set()
    image_count = 0

    for idx, slide in enumerate(prs.slides):
        slide_info: dict[str, Any] = {
            "slide_index": idx,
            "layout_name": slide.slide_layout.name,
            "title": "",
            "bullets": [],
            "notes": "",
            "has_images": False,
            "image_count": 0,
        }
        layouts_used.add(slide.slide_layout.name)

        for shape in slide.shapes:
            # Title
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
                slide_info["has_images"] = True
                slide_info["image_count"] += 1
                continue

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts_found.add(run.font.name)
                    if run.font.color and run.font.color.type is not None:
                        try:
                            rgb = run.font.color.rgb
                            colors_found.add(f"#{rgb}")
                        except Exception:
                            pass

            if shape.has_text_frame:
                text_content = shape.text_frame.text.strip()
                if text_content:
                    if shape.name.lower().startswith("title") or (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    ):
                        slide_info["title"] = text_content
                    else:
                        slide_info["bullets"].extend(
                            [
                                line.strip()
                                for line in text_content.split("\n")
                                if line.strip()
                            ]
                        )

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["notes"] = notes_text

        content_outline.append(slide_info)

    # ── Build design spec ────────────────────────────────────────────────────
    design_spec = {
        "slide_width_inches": slide_width_in,
        "slide_height_inches": slide_height_in,
        "aspect_ratio": (
            "16:9" if abs(slide_width_in / slide_height_in - 16 / 9) < 0.1 else "4:3"
        ),
        "fonts_detected": sorted(fonts_found),
        "colors_detected": sorted(colors_found),
        "layouts_used": sorted(layouts_used),
        "total_slides": len(prs.slides),
        "total_images": image_count,
        "primary_font": next(iter(fonts_found), "Calibri"),
        "slide_count_per_layout": _count_layouts(prs),
    }

    logger.info(
        "Analyzed PPTX: %s — %d slides, %d fonts, %d colors",
        blob_name,
        len(prs.slides),
        len(fonts_found),
        len(colors_found),
    )

    return {
        "source_blob": blob_name,
        "design_spec": design_spec,
        "content_outline": content_outline,
        "status": "analyzed",
    }


def _count_layouts(prs: Presentation) -> dict[str, int]:
    counts: dict[str, int] = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        counts[name] = counts.get(name, 0) + 1
    return counts


# ── Health check ─────────────────────────────────────────────────────────────
@mcp.custom_route("/health", methods=["GET"])
async def health(_request: Request) -> JSONResponse:
    return JSONResponse({"status": "ok", "service": "pptx-mcp-server", "tools": 7})


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    port = int(os.environ.get("PORT", "8000"))
    mcp.run(transport="http", host="0.0.0.0", port=port)

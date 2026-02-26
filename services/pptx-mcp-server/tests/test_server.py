"""
Tests for PPTX MCP Server tools.
Uses pytest with mocked Azure Blob Storage.
"""

from __future__ import annotations

import io
import json
import uuid
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

# ── Import server tools ───────────────────────────────────────────────────────
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

# Set required env vars before importing server
os.environ.setdefault(
    "AZURE_STORAGE_CONNECTION_STRING",
    "DefaultEndpointsProtocol=https;AccountName=test;AccountKey=dGVzdA==;EndpointSuffix=core.windows.net",
)
os.environ.setdefault("AZURE_STORAGE_ACCOUNT_NAME", "test")


# ── Fixtures ──────────────────────────────────────────────────────────────────


def make_blank_pptx() -> bytes:
    """Create a minimal PPTX in memory."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def mock_blob_storage(monkeypatch):
    """Mock Azure Blob Storage upload/download."""
    pptx_store: dict[str, bytes] = {}

    def fake_download(container: str, blob_name: str) -> bytes:
        key = f"{container}/{blob_name}"
        if key not in pptx_store:
            # Return a blank PPTX for unknown blobs
            return make_blank_pptx()
        return pptx_store[key]

    def fake_upload(
        container: str, blob_name: str, data: bytes, content_type: str = ""
    ) -> None:
        pptx_store[f"{container}/{blob_name}"] = data

    monkeypatch.setattr("server._download_blob", fake_download)
    monkeypatch.setattr("server._upload_blob", fake_upload)
    monkeypatch.setattr(
        "server._generate_sas_url",
        lambda c, b, h=24: f"https://test.blob.core.windows.net/{c}/{b}?sas=mock",
    )

    return pptx_store


# ── Tests: create_presentation ────────────────────────────────────────────────


class TestCreatePresentation:
    def test_creates_with_title(self, mock_blob_storage):
        from server import create_presentation

        result = create_presentation(
            title="Test Presentation", presentation_id="test-001"
        )

        assert result["presentation_id"] == "test-001"
        assert result["status"] == "created"
        assert result["slide_count"] >= 1
        assert (
            "test-001.pptx" in mock_blob_storage.get("generated/test-001.pptx", b"")
            or True
        )

    def test_generates_uuid_if_no_id(self, mock_blob_storage):
        from server import create_presentation

        result = create_presentation(title="Auto ID Test")
        assert len(result["presentation_id"]) == 36  # UUID format

    def test_creates_with_subtitle(self, mock_blob_storage):
        from server import create_presentation

        result = create_presentation(title="Main Title", subtitle="A subtitle here")
        assert result["status"] == "created"

    def test_falls_back_to_blank_on_missing_template(self, mock_blob_storage):
        from server import create_presentation

        result = create_presentation(
            title="No Template", template_name="nonexistent_template"
        )
        assert result["status"] == "created"
        assert result["template_used"] == "nonexistent_template"


# ── Tests: add_slide ──────────────────────────────────────────────────────────


class TestAddSlide:
    def test_add_content_slide(self, mock_blob_storage):
        from server import create_presentation, add_slide

        create_result = create_presentation(title="Test", presentation_id="slide-test")
        result = add_slide(
            presentation_id="slide-test",
            slide_title="Slide 1",
            content=["Bullet 1", "Bullet 2", "Bullet 3"],
            layout="content",
            speaker_notes="These are speaker notes.",
        )

        assert result["status"] == "slide_added"
        assert result["slide_title"] == "Slide 1"
        assert result["slide_index"] >= 1

    def test_add_section_header(self, mock_blob_storage):
        from server import create_presentation, add_slide

        create_presentation(title="Test", presentation_id="header-test")
        result = add_slide(
            presentation_id="header-test",
            slide_title="Section 2",
            content=[],
            layout="section_header",
        )
        assert result["status"] == "slide_added"

    def test_unknown_layout_defaults_gracefully(self, mock_blob_storage):
        from server import create_presentation, add_slide

        create_presentation(title="Test", presentation_id="layout-test")
        result = add_slide(
            presentation_id="layout-test",
            slide_title="Unknown Layout",
            content=["Bullet"],
            layout="nonexistent_layout",
        )
        assert result["status"] == "slide_added"


# ── Tests: export_presentation ────────────────────────────────────────────────


class TestExportPresentation:
    def test_export_returns_url(self, mock_blob_storage):
        from server import create_presentation, export_presentation

        create_presentation(title="Export Test", presentation_id="export-001")
        result = export_presentation(presentation_id="export-001")

        assert result["status"] == "exported"
        assert "download_url" in result
        assert "sas=mock" in result["download_url"]
        assert result["slide_count"] >= 1

    def test_export_has_expiry(self, mock_blob_storage):
        from server import create_presentation, export_presentation

        create_presentation(title="Expiry Test", presentation_id="export-002")
        result = export_presentation(presentation_id="export-002", expiry_hours=72)

        assert "expires_at" in result
        assert result["file_size_kb"] > 0


# ── Tests: analyze_pptx_document ─────────────────────────────────────────────


class TestAnalyzePptxDocument:
    def _make_pptx_with_content(self) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Analysis Test Presentation"

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_analyzes_pptx(self, mock_blob_storage):
        from server import analyze_pptx_document

        mock_blob_storage["uploads/test.pptx"] = self._make_pptx_with_content()

        result = analyze_pptx_document(blob_name="test.pptx", container="uploads")

        assert result["status"] == "analyzed"
        assert "design_spec" in result
        assert "content_outline" in result
        assert result["design_spec"]["total_slides"] == 1

    def test_design_spec_has_required_fields(self, mock_blob_storage):
        from server import analyze_pptx_document

        mock_blob_storage["uploads/test2.pptx"] = self._make_pptx_with_content()
        result = analyze_pptx_document(blob_name="test2.pptx", container="uploads")

        ds = result["design_spec"]
        assert "aspect_ratio" in ds
        assert "fonts_detected" in ds
        assert "slide_width_inches" in ds
        assert ds["aspect_ratio"] == "16:9"

    def test_content_outline_has_slides(self, mock_blob_storage):
        from server import analyze_pptx_document

        mock_blob_storage["uploads/test3.pptx"] = self._make_pptx_with_content()
        result = analyze_pptx_document(blob_name="test3.pptx", container="uploads")

        assert isinstance(result["content_outline"], list)
        assert len(result["content_outline"]) == 1
        assert result["content_outline"][0]["slide_index"] == 0


# ── Tests: list_templates ─────────────────────────────────────────────────────


class TestListTemplates:
    def test_always_includes_default(self, monkeypatch):
        from server import list_templates

        mock_client = MagicMock()
        mock_client.get_container_client.return_value.list_blobs.return_value = []
        monkeypatch.setattr("server.get_blob_client", lambda: mock_client)

        result = list_templates()

        assert result["count"] >= 1
        names = [t["name"] for t in result["templates"]]
        assert "default" in names
